import pandas as pd
import docx
import pptx
import PyPDF2
from langchain_experimental.agents import create_pandas_dataframe_agent
from langchain.llms import HuggingFaceEndpoint
import logging
import io

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def parse_csv(file):
    return pd.read_csv(file)

def parse_txt(file):
    return pd.DataFrame([line.decode('utf-8').strip() for line in file], columns=['Text'])

def parse_doc(file):
    doc = docx.Document(file)
    return pd.DataFrame([para.text for para in doc.paragraphs], columns=['Text'])

def parse_ppt(file):
    prs = pptx.Presentation(file)
    slides = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides.append(shape.text)
    return pd.DataFrame(slides, columns=['Text'])

def parse_pdf(file):
    reader = PyPDF2.PdfFileReader(file)
    pages = []
    for page_num in range(reader.getNumPages()):
        page = reader.getPage(page_num)
        pages.append(page.extractText())
    return pd.DataFrame(pages, columns=['Text'])

def get_answer(data, query, api_key):
    try:
        file_type = data.name.split('.')[-1].lower()
        if file_type == 'csv':
            df = parse_csv(data)
        elif file_type == 'txt':
            df = parse_txt(data)
        elif file_type in ['ppt', 'pptx']:
            df = parse_ppt(data)
        elif file_type in ['doc', 'docx']:
            df = parse_doc(data)
        elif file_type == 'pdf':
            df = parse_pdf(data)
        else:
            return f"Unsupported file type: {file_type}"

    except Exception as e:
        logger.error(f"Error parsing the file: {str(e)}")
        return f"Error parsing the file: {str(e)}"

    try:
        llm = HuggingFaceEndpoint(
            repo_id="mistralai/Mistral-7B-Instruct-v0.2",
            temperature=0.7,
            huggingfacehub_api_token=api_key
        )

        agent = create_pandas_dataframe_agent(
            llm,
            df,
            verbose=True,
            handle_parsing_errors=True,
            max_iterations=1000,
            timeout=3600
        )

        logger.info(f"Invoking agent with query: {query}")
        result = agent.invoke(query)
        logger.info(f"Agent response: {result}")
        return result

    except ValueError as ve:
        logger.error(f"ValueError: {str(ve)}")
        return f"ValueError: {str(ve)}"
    except Exception as e:
        logger.error(f"An unexpected error occurred: {str(e)}")
        return f"An unexpected error occurred: {str(e)}"
